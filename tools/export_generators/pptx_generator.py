"""PowerPoint generator — edits an uploaded template in place.

The user is required to upload a .pptx template. We open the template,
clone its first slide layout for each new slide, and inject content
while preserving the template's master, theme, fonts, colors, and any
images on the master. The result is a deck that visually matches the
template with EY-style enterprise wording.
"""
from __future__ import annotations

import copy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _clear_existing_slides(prs: Presentation) -> None:
    """Remove all slides from the template so we start with the master only."""
    xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides = list(xml_slides)
    for s in slides:
        xml_slides.remove(s)


def _pick_layout(prs: Presentation, layout_type: str):
    """Pick a slide layout by approximate name; fall back to index."""
    name_map = {
        "title": ["Title Slide", "Title", "Cover"],
        "content": ["Title and Content", "Content", "Title, Content"],
        "section": ["Section Header", "Section"],
    }
    candidates = name_map.get(layout_type, [])
    for cand in candidates:
        for layout in prs.slide_layouts:
            if cand.lower() in (layout.name or "").lower():
                return layout
    # Fallbacks
    if layout_type == "title" and len(prs.slide_layouts) > 0:
        return prs.slide_layouts[0]
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]


def _set_text(placeholder, text: str, *, font_size: int | None = None, bold: bool | None = None) -> None:
    """Populate a placeholder's first text frame, preserving template formatting."""
    if placeholder is None:
        return
    tf = placeholder.text_frame
    tf.text = text
    if font_size or bold is not None:
        for para in tf.paragraphs:
            for run in para.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold


def _add_bullets(placeholder, bullets: list[str]) -> None:
    if placeholder is None or not bullets:
        return
    tf = placeholder.text_frame
    tf.text = bullets[0]
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0


def build_pptx(
    output_path: Path,
    plan: dict[str, Any],
    template_path: Path | None = None,
) -> Path:
    """Build a .pptx using the user's uploaded template.

    Plan shape:
      {
        "title": str,
        "subtitle": str,
        "slides": [
          {"type": "title"|"section"|"content",
           "title": str, "subtitle": str?, "bullets": [str]?, "body": str?}
        ]
      }
    """
    if not template_path or not template_path.exists():
        raise ValueError("PPT template is required but was not provided.")

    prs = Presentation(str(template_path))
    _clear_existing_slides(prs)

    # ── Title slide ──
    title_layout = _pick_layout(prs, "title")
    title_slide = prs.slides.add_slide(title_layout)
    deck_title = plan.get("title") or "Presentation"
    deck_subtitle = plan.get("subtitle") or ""

    title_ph = title_slide.shapes.title
    if title_ph is not None:
        _set_text(title_ph, deck_title)

    for ph in title_slide.placeholders:
        if ph.placeholder_format.idx not in (0,) and deck_subtitle:
            _set_text(ph, deck_subtitle)
            break

    # ── Content slides ──
    content_layout = _pick_layout(prs, "content")
    section_layout = _pick_layout(prs, "section")

    for slide_def in plan.get("slides") or []:
        s_type = (slide_def.get("type") or "content").lower()
        layout = section_layout if s_type == "section" else content_layout
        slide = prs.slides.add_slide(layout)

        s_title = slide_def.get("title") or ""
        if slide.shapes.title is not None:
            _set_text(slide.shapes.title, s_title)

        bullets = slide_def.get("bullets") or []
        body = slide_def.get("body") or ""

        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0 and ph.has_text_frame:
                body_ph = ph
                break

        if body_ph is not None:
            if bullets:
                _add_bullets(body_ph, bullets)
            elif body:
                _set_text(body_ph, body)

    prs.save(output_path)
    return output_path
